"""Generate an executive PowerPoint from monthly portfolio synthesis."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from .utils import ensure_directory


class PowerPointGenerator:
    """Create a 5-7 slide client-ready executive presentation."""

    def __init__(self, output_dir: Path) -> None:
        self.output_dir = ensure_directory(output_dir)

    def build(self, monthly_summary: str, output_name: str = "monthly_report.pptx") -> Path:
        """Write a polished PowerPoint file from synthesized insights."""

        presentation = Presentation()
        self._set_theme(presentation)

        sections = self._parse_summary(monthly_summary)
        self._title_slide(presentation)
        self._bullet_slide(presentation, "Portfolio Health", sections.get("portfolio_health", []))
        self._bullet_slide(presentation, "Trend Snapshot", sections.get("trend_snapshot", []))
        self._bullet_slide(presentation, "Risk Trends", sections.get("common_risks", []) + sections.get("red_projects", []) )
        self._bullet_slide(presentation, "Delayed Milestones", sections.get("delayed_milestones", []))
        self._bullet_slide(presentation, "Emerging Risks", sections.get("emerging_risks", []))
        self._bullet_slide(presentation, "Recommendations", sections.get("recommendations", []))
        self._bullet_slide(presentation, "Next Month Outlook", sections.get("outlook", []))

        output_path = self.output_dir / output_name
        presentation.save(output_path)
        return output_path

    def _set_theme(self, presentation: Presentation) -> None:
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

    def _title_slide(self, presentation: Presentation) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = "Project Health Portfolio Review"
        subtitle = slide.placeholders[1]
        subtitle.text = "Executive synthesis for client leadership"
        subtitle.text_frame.paragraphs[0].font.size = Pt(20)

    def _bullet_slide(self, presentation: Presentation, title: str, bullets: list[str]) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        if not bullets:
            bullets = ["No material updates to report."]
        for index, bullet in enumerate(bullets[:6]):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = RGBColor(45, 55, 72)
        body.word_wrap = True

    def _parse_summary(self, monthly_summary: str) -> dict[str, list[str]]:
        sections = {
            "portfolio_health": [],
            "trend_snapshot": [],
            "common_risks": [],
            "delayed_milestones": [],
            "red_projects": [],
            "emerging_risks": [],
            "recommendations": [],
            "outlook": [],
        }
        current = None
        for raw_line in monthly_summary.splitlines():
            line = raw_line.strip()
            if not line:
                continue
            lowered = line.lower()
            if lowered.startswith("portfolio health"):
                current = "portfolio_health"
                sections[current].append(line)
            elif lowered.startswith("trend snapshot"):
                current = "trend_snapshot"
            elif lowered.startswith("common risks"):
                current = "common_risks"
            elif lowered.startswith("most delayed milestones"):
                current = "delayed_milestones"
            elif lowered.startswith("red projects"):
                current = "red_projects"
            elif lowered.startswith("projects becoming worse"):
                current = "emerging_risks"
            elif lowered.startswith("recommendations"):
                current = "recommendations"
            elif lowered.startswith("next month outlook"):
                current = "outlook"
            elif line.startswith("-") and current:
                sections[current].append(line.lstrip("- ").strip())
        if not sections["outlook"]:
            sections["outlook"] = ["Continue weekly monitoring and focus on the highest-risk workstreams."]
        return sections
